import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import seaborn as sns
import matplotlib.pyplot as plt
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches
import tempfile
import io
import numpy as np
from sklearn.ensemble import IsolationForest
from statsmodels.tsa.holtwinters import ExponentialSmoothing
import traceback

st.set_page_config(page_title="Analyse ventes contrat", layout="wide")
st.title("📊 Analyse des Ventes - Contrats et Assurances")

# ------------------------
# FONCTIONS UTILITAIRES
# ------------------------

def detect_column(possible_names, df_cols):
    for name in possible_names:
        for col in df_cols:
            if name.lower().strip() == col.lower().strip():
                return col
    return None

def standardize_columns(df, columns_mapping):
    df_std = df.copy()
    try:
        df_std["Date"] = pd.to_datetime(df[columns_mapping['date']], dayfirst=True, errors="coerce")
        df_std["Revenu"] = pd.to_numeric(df[columns_mapping['revenu']], errors="coerce")
        df_std["Marge"] = pd.to_numeric(df[columns_mapping['marge']], errors="coerce")
        df_std["Produit"] = df[columns_mapping['produit']].astype(str)
        df_std["Assureur"] = pd.to_numeric(df[columns_mapping['assureur']], errors="coerce")
        df_std["Distributeur"] = df[columns_mapping['distributeur']].astype(str)
    except Exception as e:
        st.error(f"Erreur lors de la standardisation des colonnes: {e}")
    return df_std

def create_pdf(summary_text, image_buffers):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    for line in summary_text.split("\n"):
        pdf.cell(0, 10, line, ln=True)
    for img_buf in image_buffers:
        tmp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        tmp_img.write(img_buf.getbuffer())
        tmp_img.close()
        pdf.image(tmp_img.name, w=180)
        pdf.ln(5)
    tmp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    pdf.output(tmp_pdf.name)
    return tmp_pdf.name

def create_pptx(image_buffers, summary_text):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Rapport de Ventes"
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    for line in summary_text.split("\n"):
        tf.add_paragraph().text = line
    for img_buf in image_buffers:
        slide = prs.slides.add_slide(slide_layout)
        img_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
        img_file.write(img_buf.getbuffer())
        img_file.close()
        slide.shapes.add_picture(img_file.name, Inches(1), Inches(1), width=Inches(8))
    tmp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp_pptx.name)
    return tmp_pptx.name

@st.cache_data
def convert_to_excel(df):
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="xlsxwriter") as writer:
        df.to_excel(writer, index=False, sheet_name="Analyse")
    output.seek(0)
    return output

def forecast_revenue(df_filtered):
    ts = df_filtered.groupby(df_filtered["Date"].dt.to_period("M"))["Revenu"].sum()
    ts.index = ts.index.to_timestamp()
    if len(ts) > 6:
        model = ExponentialSmoothing(ts, trend="add", seasonal=None)
        fit = model.fit()
        forecast = fit.forecast(3)
        return forecast
    return None

def detect_anomalies(df_filtered):
    if len(df_filtered) > 20:
        X = df_filtered[["Revenu", "Marge"]].fillna(0)
        clf = IsolationForest(contamination=0.05)
        preds = clf.fit_predict(X)
        anomalies = df_filtered[preds == -1]
        return anomalies
    return pd.DataFrame()

# ------------------------
# CHARGEMENT DU FICHIER
# ------------------------

uploaded_file = st.file_uploader("📂 Téléchargez votre fichier Excel ou CSV", type=["xlsx", "xls", "csv"])

if uploaded_file:
    try:
        if uploaded_file.name.endswith(".csv"):
            df = pd.read_csv(uploaded_file)
        else:
            df = pd.read_excel(uploaded_file)

        st.subheader("Aperçu du fichier chargé")
        st.dataframe(df.head())

        synonyms = {
            "date": ["date", "date début", "date de vente"],
            "revenu": ["prime total ttc", "revenu", "ca"],
            "marge": ["marge", "marge distributeur ttc"],
            "produit": ["produit", "device", "type", "categorie"],
            "assureur": ["part assureur", "part", "taux"],
            "distributeur": ["distributeur", "revendeur", "client", "point de vente"],
        }

        df_cols = list(df.columns)
        detected_cols = {
            key: detect_column(synonyms[key], df_cols) or df_cols[0]
            for key in synonyms
        }

        st.markdown("### 🔧 Confirmez ou ajustez les colonnes")
        for key in detected_cols:
            detected_cols[key] = st.selectbox(
                f"Colonne pour {key.capitalize()}",
                df_cols,
                index=df_cols.index(detected_cols[key])
            )

        df_std = standardize_columns(df, detected_cols)

        # ------------------------
        # FILTRES AVANCÉS
        # ------------------------
        st.sidebar.header("🎛️ Filtres avancés")

        min_date, max_date = df_std["Date"].min().date(), df_std["Date"].max().date()
        start_date = st.sidebar.date_input("🗓️ Date début", value=min_date, min_value=min_date, max_value=max_date)
        end_date = st.sidebar.date_input("📅 Date fin", value=max_date, min_value=min_date, max_value=max_date)

        produits_dispo = df_std["Produit"].unique().tolist()
        selected_produits = st.sidebar.multiselect("🛆 Produits à afficher", produits_dispo, default=produits_dispo)

        distributeurs_dispo = df_std["Distributeur"].unique().tolist()
        selected_distributeurs = st.sidebar.multiselect("🏪 Distributeurs à afficher", distributeurs_dispo, default=distributeurs_dispo)

        # Recherche intelligente
        search_term = st.sidebar.text_input("🔎 Rechercher un produit/distributeur")
        if search_term:
            selected_produits = [p for p in produits_dispo if search_term.lower() in p.lower()]
            selected_distributeurs = [d for d in distributeurs_dispo if search_term.lower() in d.lower()]

        df_filtered = df_std[
            (df_std["Date"].dt.date >= start_date) &
            (df_std["Date"].dt.date <= end_date) &
            (df_std["Produit"].isin(selected_produits)) &
            (df_std["Distributeur"].isin(selected_distributeurs))
        ]

        # ------------------------
        # KPIs ET COMMENTAIRES
        # ------------------------
        st.markdown("### 📌 Résumé détaillé de l'activité")
        total_revenu = df_filtered["Revenu"].sum()
        total_marge = df_filtered["Marge"].sum()
        revenu_moyen = df_filtered["Revenu"].mean()
        marge_moyenne = df_filtered["Marge"].mean()
        nb_contrats = len(df_filtered)
        date_min = df_filtered["Date"].min().date()
        date_max = df_filtered["Date"].max().date()
        nb_jours = (date_max - date_min).days + 1
        top_produit = df_filtered.groupby("Produit")["Revenu"].sum().idxmax()

        kpi = st.columns(4)
        kpi[0].metric("💰 Revenu Total", f"{total_revenu:,.2f} TND")
        kpi[1].metric("📈 Marge Totale", f"{total_marge:,.2f} TND")
        kpi[2].metric("📊 Revenu Moyen", f"{revenu_moyen:,.2f} TND")
        kpi[3].metric("🏆 Top Produit", top_produit)

        kpi2 = st.columns(4)
        kpi2[0].metric("📁 Contrats", nb_contrats)
        kpi2[1].metric("🗓️ Période", f"{date_min} ➔ {date_max}")
        kpi2[2].metric("📆 Jours couverts", nb_jours)

        # ------------------------
        # GRAPHIQUES INTERACTIFS
        # ------------------------
        st.markdown("### 📆 Évolution du revenu quotidien (interactif)")
        revenu_par_jour = df_filtered.groupby(df_filtered["Date"].dt.date)["Revenu"].sum()
        fig1 = px.line(revenu_par_jour, x=revenu_par_jour.index, y=revenu_par_jour.values,
                       labels={"x": "Date", "y": "Revenu (TND)"}, title="Revenu Quotidien")
        st.plotly_chart(fig1, use_container_width=True)

        st.markdown("### 🥇 Top 10 Produits par Revenu (interactif)")
        top_produits = df_filtered.groupby("Produit")["Revenu"].sum().nlargest(10)
        fig2 = px.bar(top_produits, x=top_produits.values, y=top_produits.index, orientation='h',
                      labels={"x": "Revenu (TND)", "y": "Produit"}, color=top_produits.values,
                      title="Top 10 Produits")
        st.plotly_chart(fig2, use_container_width=True)

        st.markdown("### 🎯 Répartition des revenus par produit (camembert interactif)")
        revenus_par_produit = df_filtered.groupby("Produit")["Revenu"].sum()
        fig3 = px.pie(values=revenus_par_produit.values, names=revenus_par_produit.index,
                      title="Part de chaque produit dans le revenu")
        st.plotly_chart(fig3, use_container_width=True)

        st.markdown("### 🔥 Heatmap Produit / Distributeur (matrice interactive)")
        pivot = df_filtered.pivot_table(index="Produit", columns="Distributeur", values="Revenu", aggfunc="sum", fill_value=0)
        fig4 = go.Figure(data=go.Heatmap(
            z=pivot.values,
            x=pivot.columns,
            y=pivot.index,
            colorscale="YlGnBu"
        ))
        fig4.update_layout(title="Heatmap Revenu Produit/Distributeur")
        st.plotly_chart(fig4, use_container_width=True)

        st.markdown("### 📅 Revenu mensuel (barres interactives)")
        df_filtered["Mois"] = df_filtered["Date"].dt.to_period("M").astype(str)
        revenu_mensuel = df_filtered.groupby("Mois")["Revenu"].sum()
        fig5 = px.bar(revenu_mensuel, x=revenu_mensuel.index, y=revenu_mensuel.values,
                      labels={"x": "Mois", "y": "Revenu (TND)"}, title="Revenu Mensuel")
        st.plotly_chart(fig5, use_container_width=True)

        st.markdown("### 📦 Dispersion des marges par produit (boxplot)")
        fig7 = px.box(df_filtered, x="Produit", y="Marge", color="Produit")
        st.plotly_chart(fig7, use_container_width=True)

        st.markdown("### 📆 Évolution mensuelle du Revenu")
        df_filtered["Mois"] = df_filtered["Date"].dt.to_period("M").astype(str)
        revenu_par_mois = df_filtered.groupby("Mois")["Revenu"].sum()
        fig5, ax5 = plt.subplots()
        revenu_par_mois.plot(kind="bar", ax=ax5, color="teal")
        ax5.set_ylabel("Revenu (TND)")
        st.pyplot(fig5)

        st.markdown("### 🏅 Top 5 Distributeurs par Revenu")
        top5_distrib = revenu_par_distrib.head(5)
        st.dataframe(top5_distrib.reset_index().rename(columns={"Distributeur": "Distributeur", "Revenu": "Revenu Total (TND)"}))

        # Prédiction
        st.markdown("### 🔮 Prévision du revenu (3 mois)")
        forecast = forecast_revenue(df_filtered)
        if forecast is not None:
            fig8 = go.Figure()
            fig8.add_trace(go.Bar(x=forecast.index, y=forecast.values, name="Prévision"))
            st.plotly_chart(fig8, use_container_width=True)
        else:
            st.info("Pas assez de données pour la prévision.")

        # Détection d'anomalies
        anomalies = detect_anomalies(df_filtered)
        if not anomalies.empty:
            st.markdown("### ⚠️ Anomalies détectées")
            st.dataframe(anomalies)

        # ------------------------
        # EXPORTS
        # ------------------------

        excel_data = convert_to_excel(df_filtered)
        st.download_button(
            "📥 Télécharger en Excel",
            data=excel_data,
            file_name="analyse_ventes.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

        # Génération d'images pour PDF/PPTX à partir de Plotly
        buffers = []
        for fig in [fig1, fig2, fig3, fig4, fig5, fig6, fig7]:
            buf = io.BytesIO()
            fig.write_image(buf, format="png")
            buf.seek(0)
            buffers.append(buf)

        summary_text = f"""
        Rapport de Ventes
        Période : {start_date} à {end_date}
        Revenu Total : {total_revenu:,.2f} TND
        Marge Totale : {total_marge:,.2f} TND
        Nombre de Contrats : {nb_contrats}
        Top Produit : {top_produit}
        """

        st.markdown("### 🧾 Générer le rapport PDF ou PowerPoint")
        colpdf, colpptx = st.columns(2)
        with colpdf:
            if st.button("📄 Télécharger le rapport PDF complet"):
                pdf_file = create_pdf(summary_text, buffers)
                with open(pdf_file, "rb") as f:
                    st.download_button("📥 Télécharger le PDF", data=f.read(), file_name="rapport_complet.pdf", mime="application/pdf")
        with colpptx:
            if st.button("📊 Télécharger le rapport PowerPoint"):
                pptx_file = create_pptx(buffers, summary_text)
                with open(pptx_file, "rb") as f:
                    st.download_button("📥 Télécharger le PPTX", data=f.read(), file_name="rapport_complet.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    except Exception as e:
        st.error(f"❌ Une erreur est survenue : {e}")
        st.exception(e)
else:
    st.info("🕐 Veuillez uploader un fichier Excel ou CSV pour commencer.")
